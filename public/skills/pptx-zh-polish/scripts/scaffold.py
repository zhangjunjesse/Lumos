"""为某一页生成 rebuild_slideN.py 骨架 (含 add_shape / add_text / set_text helpers)。

用法:
    python scaffold.py <src.pptx> <slide_idx_1based> <workdir>

生成 <workdir>/rebuild_slideN.py 和 <workdir>/styles.py (如不存在)。
直接编辑 build() 函数写布局。
"""
import os
import shutil
import sys

SKILL_ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))

TEMPLATE = '''"""Slide {N} — (在此写本页主题)。"""
import os
import sys

WORKDIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, WORKDIR)
sys.path.insert(1, {skill_root!r})

from styles import *
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt


def set_text(tf, text, font_size=FS_BODY, bold=False, color=SLATE_800,
             font=FONT_CN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
             letter_spacing=None, italic=None):
    tf.margin_left = Inches(0.04); tf.margin_right = Inches(0.04)
    tf.margin_top = Inches(0.0);   tf.margin_bottom = Inches(0.0)
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    for c in list(p._p):
        if c.tag == qn('a:r') or c.tag == qn('a:br'):
            p._p.remove(c)
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(font_size)
    r.font.bold = bold
    if italic is not None:
        r.font.italic = italic
    r.font.color.rgb = color
    rPr = r._r.get_or_add_rPr()
    for ea in rPr.findall(qn('a:ea')):
        rPr.remove(ea)
    ea = etree.SubElement(rPr, qn('a:ea'))
    ea.set('typeface', FONT_CN)
    if letter_spacing is not None:
        rPr.set('spc', str(letter_spacing))


def add_shape(slide, x, y, w, h, fill=None, line=None,
              shape=MSO_SHAPE.RECTANGLE, line_w=None):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        if line_w:
            s.line.width = Pt(line_w)
    tf = s.text_frame
    tf.margin_left = Inches(0.0); tf.margin_right = Inches(0.0)
    tf.margin_top = Inches(0.0); tf.margin_bottom = Inches(0.0)
    return s


def add_text(slide, x, y, w, h, text, **kwargs):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(tb.text_frame, text, **kwargs)
    return tb


def add_top_bar(slide, x, y, w, h, fill, adj=50000):
    """顶色条:上圆下直。"""
    s = add_shape(slide, x, y, w, h, fill=fill, shape=SHAPE_TOPBAR)
    set_topbar_adj(s, adj)
    return s


def clear_slide(slide):
    for sh in list(slide.shapes):
        sp = sh._element
        sp.getparent().remove(sp)


W, H = SLIDE_W, SLIDE_H


def build(slide):
    clear_slide(slide)
    add_shape(slide, 0, 0, W, H, fill=WHITE)

    # ===== 标题区 =====
    add_shape(slide, 0.40, 0.32, 0.62, 0.52,
              fill=SLATE_900, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text(slide, 0.40, 0.32, 0.62, 0.52,
             '{NN}', font_size=FS_H3 + 6, font=FONT_EN, color=WHITE, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 1.15, 0.30, 11.6, 0.55,
             '本页主标题',
             font_size=FS_H1, font=FONT_CN, color=SLATE_900, bold=True,
             letter_spacing=100, anchor=MSO_ANCHOR.MIDDLE)
    add_text(slide, 1.15, 0.85, 11.6, 0.32,
             '本页副标题 / 一句话点题',
             font_size=FS_H3 - 1, font=FONT_CN, color=SLATE_500,
             letter_spacing=100, anchor=MSO_ANCHOR.MIDDLE)
    add_shape(slide, 0.40, 1.28, 12.53, 0.02, fill=SLATE_200)

    # ===== 主体 =====
    # TODO: 在这里写卡片/图表/对比块
    # 常用布局见 SKILL.md "关键约定"

    # ===== 底部总结条 =====
    STRIP_Y = 5.70
    STRIP_H = 0.95
    add_shape(slide, 0.40, STRIP_Y, 12.53, STRIP_H,
              fill=SLATE_900, shape=SHAPE_CARD)
    add_shape(slide, 0.40, STRIP_Y, 0.08, STRIP_H, fill=SKY_500)
    add_text(slide, 0.70, STRIP_Y, 12.20, STRIP_H,
             '本页总结,一句话讲清结论',
             font_size=FS_H3 + 3, font=FONT_CN, color=WHITE, bold=True,
             letter_spacing=120, anchor=MSO_ANCHOR.MIDDLE)


if __name__ == '__main__':
    src = sys.argv[1]
    dst = sys.argv[2] if len(sys.argv) > 2 else src.replace('.pptx', '_new.pptx')
    p = Presentation(src)
    build(p.slides[{IDX0}])
    p.save(dst)
    print('saved', dst)
'''


def main():
    if len(sys.argv) < 4:
        print(__doc__)
        sys.exit(1)
    src = sys.argv[1]
    n = int(sys.argv[2])
    workdir = sys.argv[3]
    os.makedirs(workdir, exist_ok=True)

    # 拷贝 styles.py 作为起点(可被用户覆盖)
    styles_src = os.path.join(SKILL_ROOT, 'styles.py')
    styles_dst = os.path.join(workdir, 'styles.py')
    if not os.path.exists(styles_dst):
        shutil.copyfile(styles_src, styles_dst)
        print(f'copied -> {styles_dst}')

    out = os.path.join(workdir, f'rebuild_slide{n}.py')
    content = (TEMPLATE
               .replace('{N}', str(n))
               .replace('{NN}', f'{n:02d}')
               .replace('{IDX0}', str(n - 1))
               .replace('{skill_root!r}', repr(SKILL_ROOT)))
    with open(out, 'w') as f:
        f.write(content)
    print(f'wrote -> {out}')
    print(f'edit then run: python {out} {src} {workdir}/preview.pptx')


if __name__ == '__main__':
    main()
