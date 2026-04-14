"""批量转换形状类型(整个 deck 或单页)。

用法:
    python convert_shapes.py <src> <dst> --from roundRect,round2SameRect --to rect
    python convert_shapes.py <src> <dst> --from rect --to roundRect --slide 3
    python convert_shapes.py <src> <dst> --from round2SameRect --to rect --clear-adj
"""
import argparse

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn


def convert_slide(slide, from_set, to_prst, clear_adj=True):
    n = 0
    for sh in slide.shapes:
        for prstGeom in sh._element.iter(qn('a:prstGeom')):
            if prstGeom.get('prst') in from_set:
                prstGeom.set('prst', to_prst)
                if clear_adj:
                    avLst = prstGeom.find(qn('a:avLst'))
                    if avLst is not None:
                        for g in list(avLst):
                            avLst.remove(g)
                n += 1
    return n


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('src')
    ap.add_argument('dst')
    ap.add_argument('--from', dest='frm', required=True,
                    help='comma-separated prst names to match')
    ap.add_argument('--to', required=True,
                    help='target prst name')
    ap.add_argument('--slide', type=int, default=0,
                    help='1-based slide index; 0=all')
    ap.add_argument('--clear-adj', action='store_true', default=True)
    args = ap.parse_args()

    from_set = set(args.frm.split(','))
    p = Presentation(args.src)
    total = 0
    for i, slide in enumerate(p.slides, 1):
        if args.slide and i != args.slide:
            continue
        n = convert_slide(slide, from_set, args.to, args.clear_adj)
        print(f'  slide {i}: {n} shapes')
        total += n
    p.save(args.dst)
    print(f'total: {total}\nsaved {args.dst}')


if __name__ == '__main__':
    main()
