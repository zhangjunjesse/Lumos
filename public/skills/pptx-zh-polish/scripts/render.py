"""PIL-based approximate slide renderer (LibreOffice-free preview).

Supports: rect, roundRect, round2SameRect (top-rounded bottom-square),
ellipse, line, picture, text (with East Asian fallback).

Limitations: letter-spacing, complex groups, gradient fills, shadows ignored.
"""
import os
import sys

from lxml import etree
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

DPI = 120
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
PX_W = int(SLIDE_W_IN * DPI)
PX_H = int(SLIDE_H_IN * DPI)

FONT_PATHS = {
    'PingFang SC':     ['/System/Library/Fonts/PingFang.ttc'],
    '微软雅黑':         ['/System/Library/Fonts/PingFang.ttc'],
    'Helvetica Neue':  ['/System/Library/Fonts/Helvetica.ttc',
                        '/System/Library/Fonts/HelveticaNeue.ttc'],
    'default':         ['/System/Library/Fonts/PingFang.ttc',
                        '/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc',
                        '/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf'],
}

_font_cache = {}


def emu_to_px(v):
    return v / 914400 * DPI


def get_font(name, size_pt, bold=False):
    key = (name, size_pt, bold)
    if key in _font_cache:
        return _font_cache[key]
    size_px = int(size_pt * DPI / 72)
    paths = FONT_PATHS.get(name, FONT_PATHS['default'])
    idx = 2 if bold else 0
    font = None
    for path in paths:
        if not os.path.exists(path):
            continue
        try:
            font = ImageFont.truetype(path, size_px, index=idx)
            break
        except Exception:
            try:
                font = ImageFont.truetype(path, size_px)
                break
            except Exception:
                continue
    if font is None:
        font = ImageFont.load_default()
    _font_cache[key] = font
    return font


def parse_color(el):
    for child in el.iter():
        if child.tag == qn('a:srgbClr'):
            v = child.get('val')
            return (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
    return None


def find_spPr(sp):
    return (sp._element.find('.//' + qn('p:spPr'))
            or sp._element.find('.//' + qn('spPr')))


def shape_fill(sp):
    spPr = find_spPr(sp)
    if spPr is None:
        return None
    solid = spPr.find(qn('a:solidFill'))
    return parse_color(solid) if solid is not None else None


def shape_line(sp):
    spPr = find_spPr(sp)
    if spPr is None:
        return None, 1
    ln = spPr.find(qn('a:ln'))
    if ln is None:
        return None, 1
    solid = ln.find(qn('a:solidFill'))
    c = parse_color(solid) if solid is not None else None
    w = ln.get('w')
    w_px = max(1, int(int(w) / 12700 * DPI / 72)) if w else 1
    return c, w_px


def get_geom_info(sp):
    """Return (prst, adj1, adj2) or (None, None, None)."""
    spPr = find_spPr(sp)
    if spPr is None:
        return None, None, None
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is None:
        return None, None, None
    prst = prstGeom.get('prst')
    adj1 = adj2 = None
    avLst = prstGeom.find(qn('a:avLst'))
    if avLst is not None:
        for gd in avLst.findall(qn('a:gd')):
            fmla = gd.get('fmla', '')
            name = gd.get('name', '')
            if fmla.startswith('val '):
                v = int(fmla[4:])
                if name == 'adj1' or (name == 'adj' and adj1 is None):
                    adj1 = v
                elif name == 'adj2':
                    adj2 = v
    return prst, adj1, adj2


def draw_top_rounded_rect(draw, box, radius, fill=None, outline=None, width=1):
    """Top two corners rounded, bottom two square."""
    x0, y0, x1, y1 = box
    r = max(0, min(radius, (x1 - x0) / 2, (y1 - y0) / 2))
    if r <= 0:
        draw.rectangle(box, fill=fill, outline=outline, width=width if outline else 0)
        return
    # body
    if fill is not None:
        draw.rectangle([x0, y0 + r, x1, y1], fill=fill)
        draw.rectangle([x0 + r, y0, x1 - r, y0 + r], fill=fill)
        draw.pieslice([x0, y0, x0 + 2 * r, y0 + 2 * r], 180, 270, fill=fill)
        draw.pieslice([x1 - 2 * r, y0, x1, y0 + 2 * r], 270, 360, fill=fill)
    if outline is not None and width:
        draw.arc([x0, y0, x0 + 2 * r, y0 + 2 * r], 180, 270, fill=outline, width=width)
        draw.arc([x1 - 2 * r, y0, x1, y0 + 2 * r], 270, 360, fill=outline, width=width)
        draw.line([x0 + r, y0, x1 - r, y0], fill=outline, width=width)
        draw.line([x0, y0 + r, x0, y1], fill=outline, width=width)
        draw.line([x1, y0 + r, x1, y1], fill=outline, width=width)
        draw.line([x0, y1, x1, y1], fill=outline, width=width)


def render_text(draw, sp, x, y, w, h):
    if not sp.has_text_frame:
        return
    tf = sp.text_frame
    anchor = 'middle'
    bodyPr = sp._element.find('.//' + qn('a:bodyPr'))
    if bodyPr is not None:
        a = bodyPr.get('anchor')
        anchor = {'t': 'top', 'b': 'bottom'}.get(a, 'middle')

    lines = []
    for p in tf.paragraphs:
        text = ''.join(r.text for r in p.runs) if p.runs else p.text
        if not text:
            continue
        fsize, fbold, fcolor, fname = 14, False, (0, 0, 0), 'PingFang SC'
        if p.runs:
            r0 = p.runs[0]
            if r0.font.size is not None:
                fsize = r0.font.size.pt
            if r0.font.bold:
                fbold = True
            try:
                rgb = r0.font.color.rgb if r0.font.color else None
                if rgb:
                    fcolor = (rgb[0], rgb[1], rgb[2])
            except Exception:
                pass
            if r0.font.name:
                fname = r0.font.name
        lines.append((text, fsize, fbold, fcolor, fname, p.alignment))

    if not lines:
        return
    line_heights = [fs * DPI / 72 * 1.25 for _, fs, *_ in lines]
    total_h = sum(line_heights)
    if anchor == 'top':
        cy = y + 3
    elif anchor == 'bottom':
        cy = y + h - total_h - 3
    else:
        cy = y + (h - total_h) / 2

    for (text, fs, fb, fc, fn, align), lh in zip(lines, line_heights):
        font = get_font(fn, fs, bold=fb)
        try:
            bbox = draw.textbbox((0, 0), text, font=font)
            tw = bbox[2] - bbox[0]
        except Exception:
            tw = len(text) * fs * 0.6
        if align and 'RIGHT' in str(align):
            tx = x + w - tw - 3
        elif align and 'CENTER' in str(align):
            tx = x + (w - tw) / 2
        else:
            tx = x + 3
        draw.text((tx, cy), text, fill=fc, font=font)
        cy += lh


def render_shape(draw, img, sp):
    try:
        left, top = sp.left or 0, sp.top or 0
        width, height = sp.width or 0, sp.height or 0
    except Exception:
        return
    x, y = emu_to_px(left), emu_to_px(top)
    w, h = emu_to_px(width), emu_to_px(height)
    box = [x, y, x + w, y + h]

    st = sp.shape_type
    if st == MSO_SHAPE_TYPE.PICTURE:
        try:
            import io
            pim = Image.open(io.BytesIO(sp.image.blob)) \
                .convert('RGBA').resize((int(w), int(h)))
            img.paste(pim, (int(x), int(y)), pim)
        except Exception:
            pass
        return
    if st == MSO_SHAPE_TYPE.LINE or 'connector' in str(st).lower():
        lc, lw = shape_line(sp)
        if lc:
            draw.line([x, y, x + w, y + h], fill=lc, width=lw)
        return
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        render_text(draw, sp, x, y, w, h)
        return

    # auto shape / placeholder
    fill = shape_fill(sp)
    lc, lw = shape_line(sp)
    prst, adj1, _ = get_geom_info(sp)

    if prst in ('ellipse', 'oval'):
        if fill:
            draw.ellipse(box, fill=fill, outline=lc, width=lw if lc else 0)
        elif lc:
            draw.ellipse(box, outline=lc, width=lw)
    elif prst == 'roundRect':
        adj = adj1 if adj1 is not None else 16667
        r = min(w, h) * (adj / 100000)
        try:
            draw.rounded_rectangle(box, radius=r, fill=fill,
                                   outline=lc, width=lw if lc else 0)
        except Exception:
            draw.rectangle(box, fill=fill, outline=lc,
                           width=lw if lc else 0)
    elif prst == 'round2SameRect':
        adj = adj1 if adj1 is not None else 50000
        r = min(w, h) * (adj / 100000)
        draw_top_rounded_rect(draw, box, r, fill=fill, outline=lc,
                              width=lw if lc else 0)
    else:
        if fill:
            draw.rectangle(box, fill=fill, outline=lc,
                           width=lw if lc else 0)
        elif lc:
            draw.rectangle(box, outline=lc, width=lw)

    render_text(draw, sp, x, y, w, h)


def render_slide(slide, out_path, bg=(255, 255, 255)):
    img = Image.new('RGB', (PX_W, PX_H), bg)
    draw = ImageDraw.Draw(img, 'RGBA')
    for sp in slide.shapes:
        try:
            render_shape(draw, img, sp)
        except Exception as e:
            print(f'shape err: {e}', file=sys.stderr)
    img.save(out_path)
    return out_path


if __name__ == '__main__':
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)
    pptx_path = sys.argv[1]
    slide_idx = int(sys.argv[2]) - 1
    out = sys.argv[3] if len(sys.argv) > 3 else '/tmp/slide_render.png'
    p = Presentation(pptx_path)
    render_slide(p.slides[slide_idx], out)
    print(f'rendered {out}')
