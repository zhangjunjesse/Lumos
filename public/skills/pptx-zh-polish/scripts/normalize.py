"""在位改色 + 改形状,不重建布局。

用法:
    python normalize.py <src.pptx> <slide_idx_1based> <dst.pptx> \\
        [--colors '[["F59E0B","0EA5E9"]]'] \\
        [--top-bars]                     # 把小高度矩形改成上圆下直
        [--top-bar-max-h 0.18]           # 小高度阈值 (inch)
        [--top-bar-min-w 0.6]            # 色条最小宽度 (inch)
        [--adj 50000]                    # 上圆下直圆角量
"""
import argparse
import json
import sys

from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn


def replace_color(slide, from_hex, to_hex):
    from_hex = from_hex.upper().lstrip('#')
    to_hex = to_hex.upper().lstrip('#')
    n = 0
    for sh in slide.shapes:
        for c in sh._element.iter(qn('a:srgbClr')):
            if (c.get('val') or '').upper() == from_hex:
                c.set('val', to_hex)
                n += 1
    return n


def to_top_round_bar(shape, adj=50000):
    spPr = shape._element.find('.//' + qn('p:spPr')) or \
           shape._element.find('.//' + qn('spPr'))
    if spPr is None:
        return False
    prstGeom = spPr.find(qn('a:prstGeom'))
    if prstGeom is None:
        return False
    prstGeom.set('prst', 'round2SameRect')
    avLst = prstGeom.find(qn('a:avLst'))
    if avLst is None:
        avLst = etree.SubElement(prstGeom, qn('a:avLst'))
    for g in list(avLst):
        avLst.remove(g)
    g1 = etree.SubElement(avLst, qn('a:gd'))
    g1.set('name', 'adj1'); g1.set('fmla', f'val {adj}')
    g2 = etree.SubElement(avLst, qn('a:gd'))
    g2.set('name', 'adj2'); g2.set('fmla', 'val 0')
    return True


def make_top_bars_rounded(slide, max_h_in=0.18, min_w_in=0.6, adj=50000):
    n = 0
    for sh in slide.shapes:
        try:
            h_in = sh.height / 914400
            w_in = sh.width / 914400
        except Exception:
            continue
        if h_in > max_h_in or w_in < min_w_in:
            continue
        spPr = sh._element.find('.//' + qn('p:spPr')) or \
               sh._element.find('.//' + qn('spPr'))
        if spPr is None:
            continue
        prstGeom = spPr.find(qn('a:prstGeom'))
        if prstGeom is None or prstGeom.get('prst') != 'rect':
            continue
        if spPr.find(qn('a:solidFill')) is None:
            continue
        if to_top_round_bar(sh, adj=adj):
            n += 1
    return n


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('src')
    ap.add_argument('slide_idx', type=int, help='1-based slide index')
    ap.add_argument('dst')
    ap.add_argument('--colors', default='[]',
                    help='JSON list of [from_hex, to_hex] pairs')
    ap.add_argument('--top-bars', action='store_true',
                    help='Convert small rect bars to round2SameRect')
    ap.add_argument('--top-bar-max-h', type=float, default=0.18)
    ap.add_argument('--top-bar-min-w', type=float, default=0.6)
    ap.add_argument('--adj', type=int, default=50000)
    args = ap.parse_args()

    p = Presentation(args.src)
    slide = p.slides[args.slide_idx - 1]

    for frm, to in json.loads(args.colors):
        n = replace_color(slide, frm, to)
        print(f'  color {frm} -> {to}: {n} refs')

    if args.top_bars:
        n = make_top_bars_rounded(slide, args.top_bar_max_h,
                                  args.top_bar_min_w, args.adj)
        print(f'  top-bars rounded: {n}')

    p.save(args.dst)
    print(f'saved {args.dst}')


if __name__ == '__main__':
    main()
