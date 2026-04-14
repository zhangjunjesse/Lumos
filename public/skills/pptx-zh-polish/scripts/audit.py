"""审计 pptx:列出每页使用的色号 / 字体 / 形状 / 形状数。

用法:
    python audit.py <pptx> [slide_idx_1based]

不带 slide_idx 则审计全部。输出格式便于决定 normalize.py 的 --colors 参数。
"""
import sys
from collections import Counter

from pptx import Presentation
from pptx.oxml.ns import qn


def audit_slide(slide, idx):
    colors = Counter()
    fonts = Counter()
    shapes = Counter()
    total = 0
    for sh in slide.shapes:
        total += 1
        for c in sh._element.iter(qn('a:srgbClr')):
            v = (c.get('val') or '').upper()
            if v:
                colors[v] += 1
        for t in sh._element.iter(qn('a:rPr')):
            for tag in ('a:latin', 'a:ea', 'a:cs'):
                el = t.find(qn(tag))
                if el is not None and el.get('typeface'):
                    fonts[el.get('typeface')] += 1
        for g in sh._element.iter(qn('a:prstGeom')):
            p = g.get('prst')
            if p:
                shapes[p] += 1

    print(f'\n=== slide {idx} ({total} shapes) ===')
    print(f'colors ({len(colors)}): ' +
          ', '.join(f'{k}×{v}' for k, v in colors.most_common(20)))
    print(f'fonts:  ' +
          ', '.join(f'{k}×{v}' for k, v in fonts.most_common()))
    print(f'shapes: ' +
          ', '.join(f'{k}×{v}' for k, v in shapes.most_common()))


def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)
    p = Presentation(sys.argv[1])
    if len(sys.argv) >= 3:
        idx = int(sys.argv[2])
        audit_slide(p.slides[idx - 1], idx)
    else:
        for i, s in enumerate(p.slides, 1):
            audit_slide(s, i)


if __name__ == '__main__':
    main()
